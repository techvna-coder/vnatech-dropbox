import streamlit as st
from io import BytesIO
import time
import re
from typing import List, Dict, Any, Tuple
from collections import Counter

# OpenAI (cho embeddings)
try:
    from openai import OpenAI
    client = OpenAI(api_key=st.secrets.get("OPENAI_API_KEY", ""))
except ImportError:
    st.error("OpenAI library not installed")
    client = None

# PDF/PPTX
try:
    import PyPDF2
except ImportError:
    st.error("PyPDF2 not installed")
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    st.error("python-pptx not installed")
    Presentation = None

# Tokenizer
try:
    import tiktoken
    tokenizer = tiktoken.encoding_for_model("text-embedding-3-small")
except Exception:
    tokenizer = None
    st.warning("tiktoken not available; chunking will fall back to rough splitting.")

# ---------- Utilities ----------
def _safe_tokenize(text: str) -> List[int]:
    if tokenizer:
        return tokenizer.encode(text)
    return list(text)

def _safe_detokenize(tokens: List[int]) -> str:
    if tokenizer:
        return tokenizer.decode(tokens)
    return "".join(tokens)

def count_tokens(text: str) -> int:
    try:
        return len(_safe_tokenize(text))
    except Exception:
        return 0

def preprocess_text(text: str) -> str:
    """Chuẩn hóa text, giữ nguyên cấu trúc quan trọng"""
    try:
        text = text.replace("\f", "\n").replace("\r", "\n")
        text = text.replace("–", "-").replace("—", "-")
        text = text.replace(""", '"').replace(""", '"')
        text = text.replace("'", "'").replace("'", "'")
        
        # Chuẩn hóa whitespace nhưng giữ structure
        lines = []
        for line in text.split("\n"):
            stripped = line.strip()
            if stripped:
                lines.append(stripped)
            elif lines and lines[-1]:  # Giữ 1 dòng trống giữa paragraphs
                lines.append("")
        
        return "\n".join(lines).strip()
    except Exception:
        return text

# ---------- Advanced Metadata Extraction ----------
def _extract_key_terms(text: str, top_n: int = 20) -> List[str]:
    """Trích xuất các thuật ngữ quan trọng (technical terms, acronyms, etc.)"""
    # Tìm acronyms (2-6 chữ in hoa)
    acronyms = re.findall(r'\b[A-Z]{2,6}\b', text)
    
    # Tìm technical patterns: XX-XX-XX, ATA codes, etc.
    tech_codes = re.findall(r'\b[A-Z0-9]+-[A-Z0-9]+(?:-[A-Z0-9]+)*\b', text)
    
    # Tìm numbers with units
    measurements = re.findall(r'\b\d+(?:\.\d+)?(?:\s*(?:mm|cm|m|kg|lb|psi|bar|°C|°F|V|A|Hz|kW|hp))\b', text)
    
    # Combine và đếm frequency
    all_terms = acronyms + tech_codes + measurements
    if not all_terms:
        return []
    
    # Lấy top N terms by frequency
    term_counts = Counter(all_terms)
    return [term for term, _ in term_counts.most_common(top_n)]

def _detect_structure_elements(text: str) -> Dict[str, Any]:
    """Phát hiện các elements cấu trúc trong text"""
    lines = text.split('\n')
    
    structure = {
        "has_headers": False,
        "has_numbered_lists": False,
        "has_bullet_lists": False,
        "has_tables": False,
        "has_code_blocks": False,
        "header_count": 0,
        "list_count": 0,
    }
    
    header_pattern = r'^(?:\d+\.)*\d+\s+[A-Z]|^[A-Z][A-Z\s]+:$|^#+\s'
    numbered_list = r'^\s*\d+[\.\)]\s+'
    bullet_list = r'^\s*[-•*]\s+'
    table_separator = r'[-─│┼┌┐└┘├┤┬┴]'
    code_indicator = r'^\s*(?:def|class|function|if|for|while|return)\s+'
    
    for line in lines:
        if re.match(header_pattern, line.strip()):
            structure["has_headers"] = True
            structure["header_count"] += 1
        
        if re.match(numbered_list, line):
            structure["has_numbered_lists"] = True
            structure["list_count"] += 1
        
        if re.match(bullet_list, line):
            structure["has_bullet_lists"] = True
            structure["list_count"] += 1
        
        if re.search(table_separator, line):
            structure["has_tables"] = True
        
        if re.match(code_indicator, line):
            structure["has_code_blocks"] = True
    
    return structure

def _classify_content_type(text: str, metadata: Dict[str, Any]) -> str:
    """Phân loại loại nội dung của section"""
    text_lower = text.lower()
    
    # Technical documentation patterns
    if any(word in text_lower for word in ["procedure", "steps", "installation", "removal", "inspection"]):
        return "procedure"
    
    if any(word in text_lower for word in ["specification", "technical data", "parameters", "limits"]):
        return "specification"
    
    if any(word in text_lower for word in ["caution", "warning", "note", "important"]):
        return "safety_note"
    
    if metadata.get("has_tables"):
        return "table_data"
    
    if metadata.get("has_numbered_lists") or metadata.get("has_bullet_lists"):
        return "list_content"
    
    return "general"

# ---------- Enhanced PDF Processing ----------
def process_pdf(file_content: BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Extract text + rich metadata from PDF"""
    if PyPDF2 is None:
        raise Exception("PyPDF2 not installed")
    
    try:
        reader = PyPDF2.PdfReader(file_content)
        pages = []
        all_text_parts = []
        
        for i, page in enumerate(reader.pages, 1):
            try:
                t = page.extract_text() or ""
            except Exception as e:
                st.warning(f"Failed to extract text from page {i}: {e}")
                t = ""
            
            if t.strip():
                pages.append({
                    "number": i,
                    "text": t.strip(),
                    "char_count": len(t),
                    "word_count": len(t.split())
                })
                all_text_parts.append(f"--- Page {i} ---\n{t}")
        
        full_text = "\n".join(all_text_parts).strip()
        
        if not full_text:
            raise Exception("No text could be extracted from the PDF")
        
        # Extract key terms from full document
        key_terms = _extract_key_terms(full_text, top_n=30)
        
        # Detect structure
        structure = _detect_structure_elements(full_text)
        
        # Build comprehensive metadata
        meta = {
            "total_pages": len(reader.pages),
            "extracted_pages": len(pages),
            "has_sections": True,
            "sections": [{"type": "page", "number": p["number"], 
                         "word_count": p["word_count"]} for p in pages],
            "key_terms": key_terms,
            "structure": structure,
            "avg_words_per_page": sum(p["word_count"] for p in pages) / max(len(pages), 1),
            "total_words": sum(p["word_count"] for p in pages),
        }
        
        # Debug info
        st.caption(f"✓ PDF: {len(pages)} pages, {meta['total_words']} words, {len(key_terms)} key terms")
        
        return full_text, meta
        
    except Exception as e:
        raise Exception(f"Failed to process PDF: {str(e)}")

# ---------- Enhanced PPTX Processing ----------
def process_pptx(file_content: BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Extract text + rich metadata from PPTX"""
    if Presentation is None:
        raise Exception("python-pptx not installed")
    
    try:
        prs = Presentation(file_content)
        slides_data = []
        all_text_parts = []
        has_tables = False
        has_images = False
        
        for s_idx, slide in enumerate(prs.slides, 1):
            slide_text_parts = [f"--- Slide {s_idx} ---"]
            slide_tables = 0
            slide_images = 0
            
            # Check for title
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                if title:
                    slide_text_parts.append(f"TITLE: {title}")
            
            for shape in slide.shapes:
                try:
                    # Text content
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        # Skip if it's the title (already processed)
                        if shape != slide.shapes.title:
                            slide_text_parts.append(shape.text.strip())
                    
                    # Tables
                    if shape.shape_type == 19:  # TABLE
                        has_tables = True
                        slide_tables += 1
                        if hasattr(shape, "table"):
                            table = shape.table
                            table_rows = []
                            for row in table.rows:
                                cells = []
                                for cell in row.cells:
                                    try:
                                        cell_text = cell.text.strip()
                                        if cell_text:
                                            cells.append(cell_text)
                                    except Exception:
                                        continue
                                if cells:
                                    table_rows.append(" | ".join(cells))
                            if table_rows:
                                slide_text_parts.append("[TABLE]\n" + "\n".join(table_rows))
                    
                    # Images (count for metadata)
                    if shape.shape_type == 13:  # PICTURE
                        has_images = True
                        slide_images += 1
                        
                except Exception as e:
                    st.warning(f"Skipped shape in slide {s_idx}: {str(e)}")
                    continue
            
            slide_full_text = "\n".join(slide_text_parts)
            
            if slide_full_text.strip():
                slides_data.append({
                    "number": s_idx,
                    "title": title,
                    "text": slide_full_text,
                    "word_count": len(slide_full_text.split()),
                    "has_table": slide_tables > 0,
                    "table_count": slide_tables,
                    "image_count": slide_images,
                })
                all_text_parts.append(slide_full_text)
        
        full_text = "\n".join(all_text_parts).strip()
        
        if not full_text:
            raise Exception("No text could be extracted from the PowerPoint")
        
        # Extract key terms
        key_terms = _extract_key_terms(full_text, top_n=30)
        
        # Detect structure
        structure = _detect_structure_elements(full_text)
        structure["has_tables"] = has_tables
        
        # Build metadata
        meta = {
            "total_slides": len(list(prs.slides)),
            "extracted_slides": len(slides_data),
            "has_sections": True,
            "sections": [{"type": "slide", "number": s["number"], 
                         "title": s["title"], "word_count": s["word_count"],
                         "has_table": s["has_table"]} for s in slides_data],
            "key_terms": key_terms,
            "structure": structure,
            "has_tables": has_tables,
            "has_images": has_images,
            "avg_words_per_slide": sum(s["word_count"] for s in slides_data) / max(len(slides_data), 1),
            "total_words": sum(s["word_count"] for s in slides_data),
        }
        
        return full_text, meta
        
    except Exception as e:
        raise Exception(f"Failed to process PPTX: {str(e)}")

# ---------- Smart Semantic Chunking ----------
def _detect_natural_breaks(text: str) -> List[int]:
    """Phát hiện ranh giới tự nhiên trong text"""
    breaks = [0]
    lines = text.split('\n')
    pos = 0
    
    for i, line in enumerate(lines):
        pos += len(line) + 1
        
        if not line.strip():
            breaks.append(pos)
            continue
        
        stripped = line.strip()
        
        # Headers (dòng ngắn kết thúc bằng : hoặc all caps)
        if len(stripped) < 100:
            if stripped.endswith(':') or (stripped.isupper() and len(stripped) > 5):
                breaks.append(pos)
                continue
            
            # Section numbers: "1.2.3 Title" or "ATA 32-41-00"
            if re.match(r'^(?:\d+\.)+\d+\s+\w+|^[A-Z]{2,}\s+\d+-\d+', stripped):
                breaks.append(pos)
                continue
        
        # List items
        if re.match(r'^\s*(?:\d+[\.\)]|[-•*])\s+', line):
            breaks.append(pos)
            continue
        
        # Tables or special markers
        if any(marker in line for marker in ['---', '===', '[TABLE]', '[FIGURE]']):
            breaks.append(pos)
            continue
    
    breaks.append(len(text))
    return sorted(set(breaks))

def _find_best_split_point(text: str, target_pos: int, window: int = 200) -> int:
    """Tìm điểm cắt tốt nhất gần target_pos"""
    start = max(0, target_pos - window)
    end = min(len(text), target_pos + window)
    search_zone = text[start:end]
    
    # 1. Paragraph breaks
    para_breaks = [m.end() for m in re.finditer(r'\n\s*\n', search_zone)]
    if para_breaks:
        closest = min(para_breaks, key=lambda x: abs(x - (target_pos - start)))
        return start + closest
    
    # 2. Sentence ends
    sent_breaks = [m.end() for m in re.finditer(r'[.!?]\s+', search_zone)]
    if sent_breaks:
        closest = min(sent_breaks, key=lambda x: abs(x - (target_pos - start)))
        return start + closest
    
    # 3. Word boundaries
    word_breaks = [m.end() for m in re.finditer(r'\s+', search_zone)]
    if word_breaks:
        closest = min(word_breaks, key=lambda x: abs(x - (target_pos - start)))
        return start + closest
    
    return target_pos

def _chunk_by_semantic_boundaries(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
    """Chunk text theo semantic boundaries"""
    if not text or not text.strip():
        return []
    
    tokens = _safe_tokenize(text)
    
    if len(tokens) <= chunk_size:
        return [text.strip()] if text.strip() else []
    
    natural_breaks = _detect_natural_breaks(text)
    chunks = []
    start_char = 0
    
    while start_char < len(text):
        start_tokens = _safe_tokenize(text[:start_char])
        target_tokens = len(start_tokens) + chunk_size
        
        if target_tokens >= len(tokens):
            end_char = len(text)
        else:
            # Binary search for character position
            left, right = start_char, len(text)
            while left < right:
                mid = (left + right) // 2
                mid_tokens = len(_safe_tokenize(text[:mid]))
                if mid_tokens < target_tokens:
                    left = mid + 1
                else:
                    right = mid
            end_char = left
        
        # Find nearest natural break
        natural_candidates = [b for b in natural_breaks if start_char < b <= end_char + 100]
        if natural_candidates:
            end_char = min(natural_candidates, key=lambda x: abs(x - end_char))
        else:
            end_char = _find_best_split_point(text, end_char, window=150)
        
        chunk_text = text[start_char:end_char].strip()
        
        # Only add non-empty chunks
        if chunk_text and len(chunk_text) > 10:  # At least 10 chars
            chunks.append(chunk_text)
        
        if end_char >= len(text):
            break
        
        # Calculate overlap
        overlap_tokens = min(chunk_overlap, len(_safe_tokenize(chunk_text)) // 2)
        if overlap_tokens > 0 and chunk_text:
            chunk_tokens = _safe_tokenize(chunk_text)
            if len(chunk_tokens) >= overlap_tokens:
                overlap_text = _safe_detokenize(chunk_tokens[-overlap_tokens:])
                overlap_pos = text.rfind(overlap_text[:50], start_char, end_char)
                if overlap_pos > start_char:
                    start_char = overlap_pos
                else:
                    start_char = end_char
            else:
                start_char = end_char
        else:
            start_char = end_char
    
    # Final filter for empty chunks
    return [c for c in chunks if c and c.strip() and len(c.strip()) > 10]

def chunk_text_smart(text: str,
                     doc_metadata: Dict[str, Any],
                     chunk_size: int = 1000,
                     chunk_overlap: int = 200) -> List[Dict[str, Any]]:
    """Tách text thông minh theo sections và semantic boundaries"""
    text = preprocess_text(text)
    
    if not text or not text.strip():
        st.warning("Empty document after preprocessing")
        return []
    
    sections = []
    
    # Parse sections từ markers
    split_pages = [p for p in text.split("\n--- Page ") if p.strip()]
    split_slides = [s for s in text.split("\n--- Slide ") if s.strip()]

    if len(split_pages) > 1:  # PDF
        for blk in split_pages:
            try:
                header, *rest = blk.split("---", 1)
                number = int(header.strip())
                content = rest[0] if rest else ""
            except Exception:
                number, content = 0, blk
            if content.strip():  # Only add non-empty content
                sections.append(("page", number, content.strip()))
    elif len(split_slides) > 1:  # PPTX
        for blk in split_slides:
            try:
                header, *rest = blk.split("---", 1)
                number = int(header.strip())
                content = rest[0] if rest else ""
            except Exception:
                number, content = 0, blk
            
            # Extract title nếu có
            title = ""
            if content.startswith("TITLE:"):
                lines = content.split('\n', 1)
                title = lines[0].replace("TITLE:", "").strip()
                content = lines[1] if len(lines) > 1 else ""
            
            if content.strip():  # Only add non-empty content
                sections.append(("slide", number, content.strip(), title))
    else:
        if text.strip():  # Ensure not empty
            sections.append(("document", 0, text, ""))

    if not sections:
        st.warning("No sections found in document")
        return []

    out: List[Dict[str, Any]] = []
    total_chunks_counter = 0
    temp_chunks_per_section = []

    # Chunk từng section
    for section_data in sections:
        if len(section_data) == 3:
            stype, snum, scontent = section_data
            stitle = ""
        else:
            stype, snum, scontent, stitle = section_data
        
        # Skip empty sections
        if not scontent or not scontent.strip():
            continue
            
        smalls = _chunk_by_semantic_boundaries(scontent, chunk_size, chunk_overlap)
        
        # Filter out empty chunks
        smalls = [s for s in smalls if s and s.strip()]
        
        if smalls:
            temp_chunks_per_section.append((stype, snum, stitle, smalls))
            total_chunks_counter += len(smalls)

    if total_chunks_counter == 0:
        st.warning("No valid chunks created")
        return []

    # Gán metadata chi tiết cho từng chunk
    running_idx = 0
    for (stype, snum, stitle, smalls) in temp_chunks_per_section:
        for i, txt in enumerate(smalls):
            # Skip if somehow empty (safety check)
            if not txt or not txt.strip():
                continue
                
            # Phân loại content type
            chunk_structure = _detect_structure_elements(txt)
            content_type = _classify_content_type(txt, chunk_structure)
            
            # Extract local key terms
            local_terms = _extract_key_terms(txt, top_n=10)
            
            chunk_meta = {
                "text": txt.strip(),  # Ensure trimmed
                "chunk_index": running_idx,
                "total_chunks": total_chunks_counter,
                "section_type": stype,
                "section_number": snum,
                "section_title": stitle,
                "is_complete_section": (len(smalls) == 1),
                "token_count": count_tokens(txt),
                "word_count": len(txt.split()),
                "char_count": len(txt),
                "content_type": content_type,
                "has_headers": chunk_structure.get("has_headers", False),
                "has_lists": chunk_structure.get("has_numbered_lists", False) or chunk_structure.get("has_bullet_lists", False),
                "has_tables": chunk_structure.get("has_tables", False),
                "local_key_terms": local_terms,
            }
            out.append(chunk_meta)
            running_idx += 1

    return out

# ---------- Embeddings ----------
def get_embeddings(texts: List[str], batch_size: int = 100) -> List[List[float]]:
    """Generate embeddings with progress tracking"""
    if client is None:
        raise Exception("OpenAI client is not initialized")
    
    # Filter out empty or invalid texts
    valid_texts = []
    text_indices = []
    for idx, text in enumerate(texts):
        if text and isinstance(text, str) and text.strip():
            # Truncate very long texts (OpenAI limit ~8191 tokens)
            if len(text) > 30000:  # ~8000 tokens roughly
                text = text[:30000]
            valid_texts.append(text.strip())
            text_indices.append(idx)
    
    if not valid_texts:
        st.warning("No valid texts to embed. All texts are empty or invalid.")
        # Return zero vectors for all
        return [[0.0] * 1536 for _ in texts]
    
    # Create result array with zero vectors
    all_embeddings = [[0.0] * 1536 for _ in texts]
    total = len(valid_texts)
    
    for i in range(0, total, batch_size):
        batch = valid_texts[i:i+batch_size]
        batch_indices = text_indices[i:i+batch_size]
        
        if total > batch_size:
            progress = min(1.0, (i + len(batch)) / total)
            st.progress(progress, text=f"Generating embeddings... {i + len(batch)}/{total}")
        
        try:
            resp = client.embeddings.create(
                model="text-embedding-3-small",
                input=batch
            )
            # Place embeddings at correct indices
            for j, emb_data in enumerate(resp.data):
                original_idx = batch_indices[j]
                all_embeddings[original_idx] = emb_data.embedding
                
        except Exception as e:
            st.error(f"Embedding batch {i//batch_size + 1} failed: {e}")
            # Retry individual items in batch
            for j, text in enumerate(batch):
                original_idx = batch_indices[j]
                try:
                    # Double check text is valid before retry
                    if not text or not text.strip():
                        continue
                    resp = client.embeddings.create(
                        model="text-embedding-3-small",
                        input=[text]
                    )
                    all_embeddings[original_idx] = resp.data[0].embedding
                except Exception as retry_e:
                    st.warning(f"Failed to embed text at index {original_idx}: {retry_e}")
                    # Keep zero vector as fallback
                    pass
        
        if i + batch_size < total:
            time.sleep(0.1)
    
    return all_embeddings
